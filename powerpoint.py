from pptx import Presentation
from nltk.corpus import stopwords

# f = open("C:\\Users\\Demilade Sodimu\\Documents\\400 Level notes\\ALPHA\\CSC 415\\new\\[WEEK 1]-LECTURE-1- Introduction-2022", "rb")
# prs = Presentation(f)
prs = Presentation('C:\\Users\\Demilade Sodimu\\Documents\\400 Level notes\\ALPHA\\CSC 413\\[TOPIC 4] - Recursion')

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)

punctuations = ['(', ')', ';', ':', '[', ']', ',', '.', ' ', '', '):']
# We initialize the stopwords variable, which is a list of words like
# "The," "I," "and," etc. that don't hold much value as keywords.
stop_words = stopwords.words('english')
print(text_runs)
k2 = []
for word in text_runs:
    if " " not in word:
        k2.append(word)
        continue
    else:
        one = word.split(" ")
        for i in range(len(one)):
            k2.append(one[i])
keywords = [word for word in k2 if not word in stop_words and not word in punctuations]
print(keywords)


